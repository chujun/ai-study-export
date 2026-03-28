#!/usr/bin/env python3
"""
X理论与Y理论 演示文稿生成器
风格：简约大方，动画形象
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 颜色方案 - 简约大方
COLORS = {
    'primary': RGBColor(0x2C, 0x3E, 0x50),      # 深蓝灰 - 主色
    'secondary': RGBColor(0x34, 0x49, 0x5E),    # 次深蓝
    'accent_x': RGBColor(0xE7, 0x4C, 0x3C),     # 红色 - X理论
    'accent_y': RGBColor(0x27, 0xAE, 0x60),      # 绿色 - Y理论
    'white': RGBColor(0xFF, 0xFF, 0xFF),        # 白色
    'light_gray': RGBColor(0xF8, 0xF9, 0xFA),   # 浅灰背景
    'text_dark': RGBColor(0x2C, 0x3E, 0x50),    # 深色文字
    'text_light': RGBColor(0x95, 0xA5, 0xA6),   # 浅色文字
}

def set_shape_transparency(shape, alpha):
    """设置形状透明度"""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 添加背景色块 - 左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['accent_x']
    left_bar.line.fill.background()

    # 右侧装饰条
    right_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.7), Inches(0), Inches(0.3), Inches(7.5)
    )
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = COLORS['accent_y']
    right_bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, section_title, section_num):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景装饰 - 左侧色块
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4), Inches(7.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = COLORS['primary']
    left_bg.line.fill.background()

    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(3), Inches(5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    return slide

def add_content_slide(prs, title, content_items, x_color=None, y_color=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (文字, 颜色)
            p.text = f"• {item[0]}"
            p.font.color.rgb = item[1]
        else:
            p.text = f"• {item}"
            p.font.color.rgb = COLORS['text_dark']

        p.font.size = Pt(20)
        p.space_after = Pt(12)

    return slide

def add_comparison_slide(prs, title, x_items, y_items):
    """添加对比页 - X vs Y"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # X理论区域 - 左侧
    x_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.5)
    )
    x_bg.fill.solid()
    x_bg.fill.fore_color.rgb = COLORS['accent_x']
    x_bg.line.fill.background()

    x_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4), Inches(0.6))
    tf = x_title.text_frame
    p = tf.paragraphs[0]
    p.text = "X 理论"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    x_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4), Inches(4.3))
    tf = x_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(x_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(8)

    # Y理论区域 - 右侧
    y_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.5)
    )
    y_bg.fill.solid()
    y_bg.fill.fore_color.rgb = COLORS['accent_y']
    y_bg.line.fill.background()

    y_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4), Inches(0.6))
    tf = y_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Y 理论"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    y_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.3), Inches(4), Inches(4.3))
    tf = y_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(y_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # 表格
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).table

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    # 设置数据行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['text_dark']
            p.alignment = PP_ALIGN.CENTER
            # 交替行颜色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']

    return slide

def add_diagram_slide(prs, title, lines, diagram_type="pyramid"):
    """添加图示页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    if diagram_type == "pyramid":
        # 金字塔模型 - 使用圆角矩形模拟层次
        # 底层 - X理论（宽矩形）
        base = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.3), Inches(7), Inches(1.2)
        )
        base.fill.solid()
        base.fill.fore_color.rgb = COLORS['accent_x']
        base.line.fill.background()

        base_text = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.8))
        tf = base_text.text_frame
        p = tf.paragraphs[0]
        p.text = "底层：X理论（基础保障）"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 中层 - 过渡层
        middle = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.9), Inches(5), Inches(1.2)
        )
        middle.fill.solid()
        middle.fill.fore_color.rgb = COLORS['secondary']
        middle.line.fill.background()

        middle_text = slide.shapes.add_textbox(Inches(2.5), Inches(4.1), Inches(5), Inches(0.8))
        tf = middle_text.text_frame
        p = tf.paragraphs[0]
        p.text = "中层：规范与文化"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 顶层 - Y理论（窄矩形）
        top = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2.5), Inches(3), Inches(1.2)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = COLORS['accent_y']
        top.line.fill.background()

        top_text = slide.shapes.add_textbox(Inches(3.5), Inches(2.7), Inches(3), Inches(0.8))
        tf = top_text.text_frame
        p = tf.paragraphs[0]
        p.text = "顶层：Y理论"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 中层说明
        for i, line in enumerate(lines):
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5 + i * 0.45), Inches(8.5), Inches(0.4))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['text_dark']

    return slide

def create_presentation():
    """创建完整的演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 封面 ==========
    add_title_slide(
        prs,
        "X理论与Y理论",
        "在各个领域的典型应用"
    )

    # ========== 目录 ==========
    add_content_slide(prs, "目录", [
        "一、理论起源与命名",
        "二、古代中国的典型应用",
        "三、古代西方的典型应用",
        "四、现代西方的典型应用",
        "五、现代中国的典型应用",
        "六、丰田模式：X/Y整合典范",
        "七、军事领域的典型应用",
        "八、教育领域的典型应用",
        "九、其他领域应用",
        "十、总结与启示",
    ])

    # ========== 第一章：理论起源 ==========
    add_section_slide(prs, "理论起源与命名", 1)

    add_content_slide(prs, "X理论与Y理论的起源", [
        ("道格拉斯·麦格雷戈（Douglas McGregor）", COLORS['primary']),
        "《企业的人性面》（The Human Side of Enterprise，1960年）",
        "提出两套对立的人性假设",
        "麦格雷戈刻意使用X和Y作为中性占位符",
    ])

    add_comparison_slide(prs, "X理论 vs Y理论 核心假设",
        [
            "人天生厌恶工作",
            "必须用奖惩控制",
            "人缺乏上进心",
            "逃避责任",
            "只关心安全感"
        ],
        [
            "工作是自然的需要",
            "人有自我导向能力",
            "内在动机比外在更重要",
            "大多数愿意承担责任",
            "创造力广泛存在"
        ]
    )

    # ========== 第二章：古代中国 ==========
    add_section_slide(prs, "古代中国的典型应用", 2)

    add_comparison_slide(prs, "古代中国：商鞅变法 vs 道家儒家",
        [
            "商鞅变法：首级计功制",
            "斩首一级 = 赐田一顷",
            "连坐法：一人犯罪他人连坐",
            "轻罪重刑：弃灰于道者黥",
            "公元前293年伊阙之战斩首24万"
        ],
        [
            "庄子：逍遥游，追求精神自由",
            "陶渊明：不为五斗米折腰",
            "王阳明：致良知，知行合一",
            "海瑞：明知必死仍上疏",
            "岳家军：冻死不拆屋，饿死不掳掠"
        ]
    )

    add_content_slide(prs, "外儒内法：复合激励体系", [
        ("表层（儒家）- 道德教化、君轻民贵 -> Y理论", COLORS['accent_y']),
        ("深层（法家）- 赏罚必信、连坐监控 -> X理论", COLORS['accent_x']),
        "汉宣帝: 汉家自有制度，本以霸王道杂之",
    ])

    # ========== 第三章：古代西方 ==========
    add_section_slide(prs, "古代西方的典型应用", 3)

    add_comparison_slide(prs, "古代西方：斯巴达 vs 苏格拉底",
        [
            "斯巴达军事共同体",
            "7岁离开家庭公共抚养",
            "希洛人：永久的'大棒'",
            "克里普提：刺杀希洛人",
            "军事优势持续约200年"
        ],
        [
            "苏格拉底：产婆术，唤醒内在认知",
            "第欧根尼：别挡住我的阳光",
            "马可·奥勒留：皇帝也是哲学家",
            "亚里士多德：幸福是灵魂合德性的活动",
            "斯多葛：内心安宁是终极追求"
        ]
    )

    add_content_slide(prs, "罗马军团：X/Y混合体系", [
        ("胡萝卜 - 军饷、战利品25%、退役金、授田", COLORS['accent_y']),
        ("大棒 - 十抽一杀（哗变部队10%死亡率）", COLORS['accent_x']),
        ("Y理论元素 - 鹰旗的神圣象征，为罗马而战的公民精神", COLORS['accent_y']),
        "波利比乌斯：罗马军团优势在于严格的纪律和激励机制",
    ])

    # ========== 第四章：现代西方 ==========
    add_section_slide(prs, "现代西方的典型应用", 4)

    add_comparison_slide(prs, "现代西方：泰勒制 vs 谷歌",
        [
            "泰勒制：科学管理四大原则",
            "动作研究、时间研究",
            "差异化计件工资",
            "福特T型车：装配93分钟",
            "亚马逊：监控手环、精确到秒"
        ],
        [
            "谷歌：心理安全是最高因素",
            "亚里士多德项目：180个团队研究",
            "20%自由时间：Gmail等诞生",
            "弹性工作制、扁平化结构",
            "微软纳德拉：成长型思维转型"
        ]
    )

    add_content_slide(prs, "现代西方成果数据", [
        ("谷歌20%时间成果 - Gmail、Google News、AdSense", COLORS['accent_y']),
        ("微软转型 - 市值3000亿->2.8万亿美元", COLORS['accent_y']),
        ("员工敬业度 - 46%->91%", COLORS['accent_y']),
        ("亚马逊争议 - 监控手环、裁员算法", COLORS['accent_x']),
    ])

    # ========== 第五章：现代中国 ==========
    add_section_slide(prs, "现代中国的典型应用", 5)

    add_comparison_slide(prs, "现代中国：计划经济 vs 改革开放",
        [
            "1949-1978年：极端X理论",
            "统包统配、终身就业",
            "政治审查、户籍控制",
            "大跃进激励失败",
            "信息扭曲导致结构崩溃"
        ],
        [
            "1978年包产到户",
            "分田到户，剩余归己",
            "粮食3亿吨→4亿吨",
            "华为：虚拟股+狼性文化",
            "字节：Context not Control"
        ]
    )

    add_content_slide(prs, "中国互联网企业的Y理论探索", [
        ("海底捞 - 员工授权、冻死不拆屋式服务", COLORS['accent_y']),
        ("员工流失率 - 10%（行业60%+）", COLORS['accent_y']),
        ("字节跳动 - OKR透明、内部论坛", COLORS['accent_y']),
        ("困境 - 996文化暴露口号与实操分裂", COLORS['accent_x']),
    ])

    # ========== 第六章：丰田模式 ==========
    add_section_slide(prs, "丰田模式：X/Y整合典范", 6)

    add_content_slide(prs, "丰田生产方式（TPS）", [
        ("支柱一：JIT准时化 - 拉动式生产、单件流", COLORS['primary']),
        ("支柱二：自働化 - 人字旁的自动化、停线权", COLORS['primary']),
        ("地基 - 尊重人性、持续改善", COLORS['primary']),
        "大野耐一：让异常停下来，是改善的前提",
    ])

    add_diagram_slide(prs, "丰田模式的整合智慧", [
        "底层：X理论 = 标准作业、流程控制、绩效考核",
        "顶层：Y理论 = 持续改善、停线权、看板拉动",
        "结论：X理论是骨骼，Y理论是灵魂",
    ], "pyramid")

    add_content_slide(prs, "丰田改善提案制度（2022）", [
        ("日本本土 - 每员工45条/年，85%实施", COLORS['accent_y']),
        ("北美 - 12条/年，70%实施", COLORS['accent_y']),
        ("中国 - 8条/年，60%实施", COLORS['accent_y']),
        ("奖励结构 - 以认可和成长为主，金钱为辅", COLORS['primary']),
    ])

    # ========== 第七章：军事 ==========
    add_section_slide(prs, "军事领域的典型应用", 7)

    add_comparison_slide(prs, "军事激励：X vs Y",
        [
            "秦军：首级计功+连坐法",
            "斯巴达：军事共同体+希洛人",
            "拿破仑：征俄失败（无利可图）",
            "X理论有效条件：短期、低烈度"
        ],
        [
            "岳家军：忠义驱动，以少胜多",
            "志愿军：炒面+雪，坚守43天",
            "以色列：全民兵役，使命感召",
            "Y理论有效条件：长期、高烈度"
        ]
    )

    add_content_slide(prs, "军事激励的层次模型", [
        ("Y理论层 - 信仰/意识形态、荣誉感、战友生死依托", COLORS['accent_y']),
        ("X理论层 - 赏罚制度、训练习惯、物质保障", COLORS['accent_x']),
        ("结论 - Y理论决定上限，X理论决定下限", COLORS['primary']),
        "失去Y = 乌合之众；失去X = 行尸走肉",
    ])

    # ========== 第八章：教育 ==========
    add_section_slide(prs, "教育领域的典型应用", 8)

    add_comparison_slide(prs, "教育：应试教育 vs 蒙特梭利",
        [
            "题海战术：重复强化",
            "排名制度：外部竞争",
            "惩罚性作业：错误惩罚",
            "升学率考核：学校被指标驱动",
            "PISA：中国数学第一、幸福感倒数"
        ],
        [
            "儿童自主选择学习内容",
            "无考试、无作业",
            "内在满足感驱动",
            "蒙特梭利毕业生：创造力更强",
            "芬兰：玩耍优先、晚入学7岁"
        ]
    )

    # ========== 第九章：其他领域 ==========
    add_section_slide(prs, "其他领域的典型应用", 9)

    add_table_slide(prs, "各领域X/Y理论应用一览",
        ["领域", "X理论代表", "Y理论代表"],
        [
            ["体育", "兴奋剂、军事化训练", "科比凌晨四点训练"],
            ["宗教", "赎罪券、异端裁判所", "因信称义、禅宗顿悟"],
            ["政府", "法治、惩罚机制", "公民责任、志愿服务"],
            ["非营利", "福利依赖", "特蕾莎修女使命感召"],
            ["家庭", "体罚、物质奖励", "正面管教、无条件接纳"],
            ["创业", "对赌协议、淘汰压力", "使命驱动、愿景驱动"],
        ]
    )

    # ========== 第十章：总结 ==========
    add_section_slide(prs, "总结与启示", 10)

    add_comparison_slide(prs, "X理论与Y理论对比",
        [
            "人性假设：消极",
            "核心驱动力：外在奖惩",
            "管理重心：控制行为",
            "适用场景：简单重复劳动",
            "领导风格：命令式",
            "监督程度：严密监控"
        ],
        [
            "人性假设：积极",
            "核心驱动力：内在动机",
            "管理重心：释放潜能",
            "适用场景：知识创新工作",
            "领导风格：授权式",
            "监督程度：高度信任"
        ]
    )

    add_diagram_slide(prs, "最佳实践：混合模式", [
        "底层：X理论 — 合理薪酬、清晰职责、必要纪律",
        "顶层：Y理论 — 使命愿景、成长机会、认可表彰",
    ], "pyramid")

    add_content_slide(prs, "核心启示", [
        ("1. 非此即彼是错误的 - 成熟组织需要X/Y动态平衡", COLORS['primary']),
        ("2. X理论是必要条件 - 防止混乱、维持秩序", COLORS['accent_x']),
        ("3. Y理论是充分条件 - 激发超越、持续创新", COLORS['accent_y']),
        ("4. 情境决定配比 - 初创期Y为主，成熟期X增加", COLORS['primary']),
        ("5. Y理论需要前提 - 高素质员工+心理安全文化", COLORS['accent_y']),
    ])

    # ========== 结束页 ==========
    add_title_slide(prs, "谢谢观看", "X理论与Y理论 · 在各个领域的典型应用")

    # 保存文件
    output_path = "/root/ai/claudecode/study-export/docs/X理论Y理论应用.pptx"
    prs.save(output_path)
    print(f"演示文稿已保存至: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
